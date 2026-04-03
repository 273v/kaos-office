#!/usr/bin/env python3
"""Generate comprehensive PPTX test fixtures using python-pptx.

Creates files that exercise every code path in the PPTX reader.
Run with: uv run python tests/generate_battle_test_pptx.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

FIXTURES_DIR = Path(__file__).parent / "fixtures" / "pptx" / "battle"


def make_rich_text_pptx():
    """PPTX with comprehensive text formatting on a single slide."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Title text box
    from pptx.util import Inches

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Rich Text Formatting Test"
    run.font.bold = True
    run.font.size = Pt(28)

    # Bold text
    p2 = tf.add_paragraph()
    run = p2.add_run()
    run.text = "This is bold text"
    run.font.bold = True

    # Italic text
    p3 = tf.add_paragraph()
    run = p3.add_run()
    run.text = "This is italic text"
    run.font.italic = True

    # Bold + italic
    p4 = tf.add_paragraph()
    run = p4.add_run()
    run.text = "This is bold and italic"
    run.font.bold = True
    run.font.italic = True

    # Multiple runs in one paragraph
    p5 = tf.add_paragraph()
    run1 = p5.add_run()
    run1.text = "Normal, "
    run2 = p5.add_run()
    run2.text = "bold"
    run2.font.bold = True
    run3 = p5.add_run()
    run3.text = ", and "
    run4 = p5.add_run()
    run4.text = "italic"
    run4.font.italic = True
    run5 = p5.add_run()
    run5.text = " in one paragraph."

    # Hyperlink
    p6 = tf.add_paragraph()
    run = p6.add_run()
    run.text = "Click here for link"
    hlink = run.hyperlink
    hlink.address = "https://273ventures.com"

    # Second text box with different position
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
    tf2 = txBox2.text_frame
    p = tf2.paragraphs[0]
    p.text = "Unicode: café résumé naïve 日本語 中文 العربية"

    p2 = tf2.add_paragraph()
    p2.text = "Special chars: <>&\"' © ® ™ — – … € £ ¥"

    p3 = tf2.add_paragraph()
    p3.text = ""  # Empty paragraph

    p4 = tf2.add_paragraph()
    p4.text = "After empty paragraph"

    path = FIXTURES_DIR / "rich-text.pptx"
    prs.save(str(path))
    print(f"  Created {path.name}")


def make_bullets_pptx():
    """PPTX with comprehensive bullet/list scenarios."""
    prs = Presentation()

    # Slide 1: Simple bullet list (via XML manipulation)
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    title = slide.shapes.title
    title.text = "Bullet Lists"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    # Level 0 bullets
    p = tf.paragraphs[0]
    p.text = "First bullet"
    p.level = 0

    p2 = tf.add_paragraph()
    p2.text = "Second bullet"
    p2.level = 0

    # Level 1 (nested)
    p3 = tf.add_paragraph()
    p3.text = "Nested item one"
    p3.level = 1

    p4 = tf.add_paragraph()
    p4.text = "Nested item two"
    p4.level = 1

    # Back to level 0
    p5 = tf.add_paragraph()
    p5.text = "Third bullet"
    p5.level = 0

    # Level 2 (deep)
    p6 = tf.add_paragraph()
    p6.text = "Deep nested"
    p6.level = 2

    # Slide 2: Numbered list using XML manipulation
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    title2.text = "Numbered Lists"

    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.clear()

    # Add numbered bullets via XML
    for i, text in enumerate(["Step one", "Step two", "Step three"]):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = text
        p.level = 0
        # Set auto-numbering
        pPr = p._p.get_or_add_pPr()
        pPr.set("lvl", "0")
        buAutoNum = pPr.makeelement(qn("a:buAutoNum"), {})
        buAutoNum.set("type", "arabicPeriod")
        pPr.append(buAutoNum)

    path = FIXTURES_DIR / "bullets.pptx"
    prs.save(str(path))
    print(f"  Created {path.name}")


def make_tables_pptx():
    """PPTX with various table configurations."""
    prs = Presentation()

    # Slide 1: Simple table
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    rows, cols = 4, 3
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(3)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Header row
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Role"
    table.cell(0, 2).text = "Department"

    # Data rows
    data = [
        ["Alice", "Engineer", "Engineering"],
        ["Bob", "Manager", "Product"],
        ["Carol", "Designer", "Design"],
    ]
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # Slide 2: Table with merged cells (via XML)
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    table_shape2 = slide2.shapes.add_table(3, 3, left, top, width, height)
    table2 = table_shape2.table

    table2.cell(0, 0).text = "Merged Header"
    table2.cell(0, 1).text = ""
    table2.cell(0, 2).text = "Single"
    table2.cell(1, 0).text = "A"
    table2.cell(1, 1).text = "B"
    table2.cell(1, 2).text = "C"
    table2.cell(2, 0).text = "D"
    table2.cell(2, 1).text = "E"
    table2.cell(2, 2).text = "F"

    # Merge cells via XML
    tc0 = table2.cell(0, 0)._tc
    tc0.set("gridSpan", "2")
    tc1 = table2.cell(0, 1)._tc
    tc1.set("hMerge", "1")

    # Slide 3: Empty table
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    table_shape3 = slide3.shapes.add_table(2, 2, left, top, width, height)
    # Leave all cells empty

    path = FIXTURES_DIR / "tables.pptx"
    prs.save(str(path))
    print(f"  Created {path.name}")


def make_charts_pptx():
    """PPTX with various chart types."""
    from pptx.chart.data import CategoryChartData

    prs = Presentation()

    # Slide 1: Bar chart with data
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
    chart_data.add_series("Revenue", (100, 200, 150, 300))
    chart_data.add_series("Costs", (80, 120, 110, 180))

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(8),
        Inches(5),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Quarterly Performance"

    # Slide 2: Pie chart
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    pie_data = CategoryChartData()
    pie_data.categories = ["Engineering", "Sales", "Marketing", "Operations"]
    pie_data.add_series("Headcount", (45, 30, 15, 10))

    chart_frame2 = slide2.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(1),
        Inches(1),
        Inches(8),
        Inches(5),
        pie_data,
    )
    chart2 = chart_frame2.chart
    chart2.has_title = True
    chart2.chart_title.text_frame.text = "Department Distribution"

    # Slide 3: Line chart
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    line_data = CategoryChartData()
    line_data.categories = ["Jan", "Feb", "Mar", "Apr", "May"]
    line_data.add_series("2024", (10, 20, 15, 25, 30))
    line_data.add_series("2025", (12, 18, 22, 28, 35))

    chart_frame3 = slide3.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(1),
        Inches(1),
        Inches(8),
        Inches(5),
        line_data,
    )

    path = FIXTURES_DIR / "charts.pptx"
    prs.save(str(path))
    print(f"  Created {path.name}")


def make_images_pptx():
    """PPTX with images and alt text."""
    import io

    from PIL import Image

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Create a simple PNG in memory
    img = Image.new("RGB", (200, 100), color="red")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)

    # Save to temp file (python-pptx needs a path)
    img_path = FIXTURES_DIR / "_temp_image.png"
    img_path.write_bytes(buf.getvalue())

    pic = slide.shapes.add_picture(str(img_path), Inches(1), Inches(1), Inches(4), Inches(2))

    # Set alt text via XML
    nvPicPr = pic._element.find(qn("p:nvPicPr"))
    if nvPicPr is not None:
        cNvPr = nvPicPr.find(qn("p:cNvPr"))
        if cNvPr is not None:
            cNvPr.set("descr", "A red rectangle test image")

    # Add a second image
    img2 = Image.new("RGB", (100, 100), color="blue")
    buf2 = io.BytesIO()
    img2.save(buf2, format="JPEG")
    buf2.seek(0)
    img2_path = FIXTURES_DIR / "_temp_image2.jpg"
    img2_path.write_bytes(buf2.getvalue())

    slide.shapes.add_picture(str(img2_path), Inches(6), Inches(1), Inches(2), Inches(2))

    # Cleanup temp files
    img_path.unlink()
    img2_path.unlink()

    path = FIXTURES_DIR / "images.pptx"
    prs.save(str(path))
    print(f"  Created {path.name}")


def make_notes_pptx():
    """PPTX with speaker notes on various slides."""
    prs = Presentation()

    # Slide 1: Has notes
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Slide With Notes"
    slide.placeholders[1].text = "Main content here"

    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = (
        "These are the speaker notes for slide 1.\nWith multiple lines.\nAnd a third line."
    )

    # Slide 2: No notes
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Slide Without Notes"
    slide2.placeholders[1].text = "No notes here"

    # Slide 3: Empty notes
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Slide With Empty Notes"
    notes_slide3 = slide3.notes_slide
    notes_slide3.notes_text_frame.text = ""

    path = FIXTURES_DIR / "notes.pptx"
    prs.save(str(path))
    print(f"  Created {path.name}")


def make_groups_pptx():
    """PPTX with group shapes containing text and images."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Add multiple text boxes that overlap (testing spatial ordering)
    shapes = []
    for i in range(5):
        txBox = slide.shapes.add_textbox(
            Inches(0.5 + i * 0.5),
            Inches(4 - i * 0.5),  # Decreasing Y = higher position
            Inches(3),
            Inches(0.5),
        )
        txBox.text_frame.text = f"Shape {i + 1} (y={4 - i * 0.5:.1f})"
        shapes.append(txBox)

    # Also add title
    txBox_title = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.8))
    txBox_title.text_frame.text = "Spatial Ordering Test"
    run = txBox_title.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(24)
    run.font.bold = True

    path = FIXTURES_DIR / "spatial-ordering.pptx"
    prs.save(str(path))
    print(f"  Created {path.name}")


def make_many_slides_pptx():
    """PPTX with many slides for performance testing."""
    prs = Presentation()

    for i in range(50):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Slide {i + 1}: Performance Test"
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        # Add multiple paragraphs
        for j in range(10):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = (
                f"Bullet point {j + 1} on slide {i + 1} with some filler text to test throughput"
            )
            p.level = j % 3  # Varying indent levels

    path = FIXTURES_DIR / "50-slides.pptx"
    prs.save(str(path))
    print(f"  Created {path.name}")


def make_mixed_content_pptx():
    """PPTX with mixed content types on a single slide."""
    from pptx.chart.data import CategoryChartData

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Mixed Content Slide"
    run.font.size = Pt(28)
    run.font.bold = True

    # Text with bullets (via XML)
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(2))
    tf2 = txBox2.text_frame
    for i, text in enumerate(["Key finding one", "Key finding two", "Key finding three"]):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = text
        # Add bullet character via XML
        pPr = p._p.get_or_add_pPr()
        pPr.set("lvl", "0")
        buChar = pPr.makeelement(qn("a:buChar"), {})
        buChar.set("char", "•")
        pPr.append(buChar)

    # Table
    table_shape = slide.shapes.add_table(3, 2, Inches(5.5), Inches(1.2), Inches(4), Inches(1.5))
    table = table_shape.table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Revenue"
    table.cell(1, 1).text = "$1.2M"
    table.cell(2, 0).text = "Growth"
    table.cell(2, 1).text = "15%"

    # Chart
    chart_data = CategoryChartData()
    chart_data.categories = ["2023", "2024", "2025"]
    chart_data.add_series("Actual", (80, 100, 120))

    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5),
        Inches(3.5),
        Inches(4.5),
        Inches(3),
        chart_data,
    )

    path = FIXTURES_DIR / "mixed-content.pptx"
    prs.save(str(path))
    print(f"  Created {path.name}")


def make_metadata_pptx():
    """PPTX with rich metadata."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Metadata Test"
    slide.placeholders[1].text = "Testing metadata extraction"

    # Set metadata
    prs.core_properties.title = "Battle Test Presentation"
    prs.core_properties.author = "Test Author"
    prs.core_properties.subject = "Battle Testing"
    prs.core_properties.comments = "Generated for testing"
    prs.core_properties.category = "Test"

    path = FIXTURES_DIR / "metadata.pptx"
    prs.save(str(path))
    print(f"  Created {path.name}")


def make_edge_cases_pptx():
    """PPTX with edge case content."""
    prs = Presentation()

    # Slide 1: Very long text
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
    tf = txBox.text_frame
    tf.word_wrap = True
    long_text = "This is a very long paragraph. " * 100
    tf.paragraphs[0].text = long_text

    # Slide 2: Single character
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(1), Inches(1))
    txBox2.text_frame.text = "X"

    # Slide 3: Only whitespace
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3 = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    txBox3.text_frame.text = "   \n\t\n   "

    # Slide 4: Multiple shapes, same position (overlap)
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    for i in range(3):
        txBox = slide4.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        txBox.text_frame.text = f"Overlapping shape {i + 1}"

    path = FIXTURES_DIR / "edge-cases.pptx"
    prs.save(str(path))
    print(f"  Created {path.name}")


def main():
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    print("Generating battle test PPTX fixtures...")

    make_rich_text_pptx()
    make_bullets_pptx()
    make_tables_pptx()
    make_charts_pptx()
    make_images_pptx()
    make_notes_pptx()
    make_groups_pptx()
    make_many_slides_pptx()
    make_mixed_content_pptx()
    make_metadata_pptx()
    make_edge_cases_pptx()

    print(f"\nGenerated {len(list(FIXTURES_DIR.glob('*.pptx')))} PPTX fixtures in {FIXTURES_DIR}")


if __name__ == "__main__":
    main()
