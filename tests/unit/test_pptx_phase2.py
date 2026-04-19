"""Unit tests for PPTX Phase 2: images, table merging, speaker notes.

Phase 1 covered basic block serialization. Phase 2 adds:
- Figure/Image → slide.shapes.add_picture
- Table cell merging via gridSpan / hMerge / rowSpan / vMerge
- Speaker notes with inline formatting (not just plain text)
"""

from __future__ import annotations

from pathlib import Path

import pytest
from kaos_content.model.attr import Attr
from kaos_content.model.blocks import Div, Figure, Heading, Paragraph, Table
from kaos_content.model.document import ContentDocument, DocumentMetadata
from kaos_content.model.inlines import Image, Text
from kaos_content.model.table import Cell, Row, TableSection
from pptx import Presentation as PptxPresentation

from kaos_office.pptx.reader import parse_pptx
from kaos_office.pptx.writer import write_pptx


@pytest.fixture
def red_png(tmp_path: Path) -> Path:
    """Create a 100x50 red PNG for image tests."""
    from PIL import Image as PILImage

    path = tmp_path / "red.png"
    PILImage.new("RGB", (100, 50), color="red").save(path)
    return path


# ---------------------------------------------------------------------------
# Images
# ---------------------------------------------------------------------------


class TestImages:
    def test_figure_with_image_becomes_picture_shape(self, tmp_path: Path, red_png: Path) -> None:
        doc = ContentDocument(
            metadata=DocumentMetadata(title=""),
            body=(
                Heading(depth=1, children=(Text(value="Image Slide"),)),
                Figure(
                    children=(Paragraph(children=(Image(src=str(red_png), alt="Red test image"),)),)
                ),
            ),
        )
        out = tmp_path / "img.pptx"
        write_pptx(doc, out)

        prs = PptxPresentation(str(out))
        assert len(prs.slides) == 1
        slide = prs.slides[0]
        # shape_type 13 == PICTURE
        pictures = [s for s in slide.shapes if s.shape_type == 13]
        assert len(pictures) == 1

    def test_image_alt_text_set(self, tmp_path: Path, red_png: Path) -> None:
        doc = ContentDocument(
            metadata=DocumentMetadata(title=""),
            body=(
                Heading(depth=1, children=(Text(value="Slide"),)),
                Figure(
                    children=(Paragraph(children=(Image(src=str(red_png), alt="Red rectangle"),)),)
                ),
            ),
        )
        out = tmp_path / "alt.pptx"
        write_pptx(doc, out)

        prs = PptxPresentation(str(out))
        slide = prs.slides[0]
        pictures = [s for s in slide.shapes if s.shape_type == 13]
        assert pictures
        # Descriptor attribute in the picture's cNvPr element
        el = pictures[0]._element
        descr_cells = el.findall(".//{*}cNvPr")
        assert any(c.get("descr") == "Red rectangle" for c in descr_cells)

    def test_missing_image_falls_back_to_alt_text(self, tmp_path: Path) -> None:
        """Unresolvable image src (e.g. pptx://) should not crash."""
        doc = ContentDocument(
            metadata=DocumentMetadata(title=""),
            body=(
                Heading(depth=1, children=(Text(value="Slide"),)),
                Figure(
                    children=(
                        Paragraph(
                            children=(Image(src="pptx://nonexistent.png", alt="Fallback text"),)
                        ),
                    )
                ),
            ),
        )
        out = tmp_path / "missing.pptx"
        write_pptx(doc, out)
        # Just ensure no crash — alt text drops into a text box
        prs = PptxPresentation(str(out))
        assert len(prs.slides) == 1

    def test_image_roundtrip_through_reader(self, tmp_path: Path, red_png: Path) -> None:
        """Image written → reader produces Image(src='pptx://...', alt=...)."""
        doc = ContentDocument(
            metadata=DocumentMetadata(title=""),
            body=(
                Heading(depth=1, children=(Text(value="Pic"),)),
                Figure(children=(Paragraph(children=(Image(src=str(red_png), alt="R"),)),)),
            ),
        )
        out = tmp_path / "rt.pptx"
        write_pptx(doc, out)
        doc2 = parse_pptx(out)

        # Walk and find at least one Image node
        from kaos_content.traversal import walk

        images = [n for block in doc2.body for n in walk(block) if isinstance(n, Image)]
        assert images, "no image survived round-trip"
        # Alt text preserved
        assert any("R" in (img.alt or "") for img in images)


# ---------------------------------------------------------------------------
# Table cell merging
# ---------------------------------------------------------------------------


class TestTableMerging:
    def test_horizontal_merge_emits_gridSpan_hMerge(self, tmp_path: Path) -> None:
        doc = ContentDocument(
            metadata=DocumentMetadata(title=""),
            body=(
                Table(
                    head=TableSection(
                        rows=(
                            Row(
                                cells=(
                                    Cell(
                                        col_span=2,
                                        content=(
                                            Paragraph(children=(Text(value="Merged Header"),)),
                                        ),
                                    ),
                                    Cell(content=(Paragraph(children=(Text(value="Solo"),)),)),
                                )
                            ),
                        )
                    ),
                    bodies=(
                        TableSection(
                            rows=(
                                Row(
                                    cells=(
                                        Cell(content=(Paragraph(children=(Text(value="a"),)),)),
                                        Cell(content=(Paragraph(children=(Text(value="b"),)),)),
                                        Cell(content=(Paragraph(children=(Text(value="c"),)),)),
                                    )
                                ),
                            )
                        ),
                    ),
                ),
            ),
        )
        out = tmp_path / "merge.pptx"
        write_pptx(doc, out)

        prs = PptxPresentation(str(out))
        slide = prs.slides[0]
        tables = [s for s in slide.shapes if s.has_table]
        assert tables
        tbl = tables[0].table
        # Row 0: cell 0 gridSpan=2, cell 1 hMerge=1, cell 2 no merge
        assert tbl.cell(0, 0)._tc.get("gridSpan") == "2"
        assert tbl.cell(0, 1)._tc.get("hMerge") == "1"
        assert tbl.cell(0, 2)._tc.get("gridSpan") is None

    def test_vertical_merge_emits_rowSpan_vMerge(self, tmp_path: Path) -> None:
        doc = ContentDocument(
            metadata=DocumentMetadata(title=""),
            body=(
                Table(
                    head=TableSection(
                        rows=(
                            Row(
                                cells=(
                                    Cell(
                                        row_span=2,
                                        content=(Paragraph(children=(Text(value="Tall Cell"),)),),
                                    ),
                                    Cell(content=(Paragraph(children=(Text(value="A1"),)),)),
                                )
                            ),
                            Row(
                                cells=(
                                    Cell(content=(Paragraph(children=(Text(value="B1"),)),)),
                                    Cell(content=(Paragraph(children=(Text(value="B2"),)),)),
                                )
                            ),
                        )
                    ),
                ),
            ),
        )
        out = tmp_path / "vmerge.pptx"
        write_pptx(doc, out)

        prs = PptxPresentation(str(out))
        tbl = prs.slides[0].shapes[0].table
        assert tbl.cell(0, 0)._tc.get("rowSpan") == "2"
        assert tbl.cell(1, 0)._tc.get("vMerge") == "1"


# ---------------------------------------------------------------------------
# Speaker notes
# ---------------------------------------------------------------------------


class TestSpeakerNotes:
    def test_speaker_notes_div_attached_to_slide(self, tmp_path: Path) -> None:
        doc = ContentDocument(
            metadata=DocumentMetadata(title=""),
            body=(
                Heading(depth=1, children=(Text(value="Slide"),)),
                Paragraph(children=(Text(value="Main text."),)),
                Div(
                    attr=Attr(classes=("speaker-notes",)),
                    children=(Paragraph(children=(Text(value="Remember to emphasize X."),)),),
                ),
            ),
        )
        out = tmp_path / "notes.pptx"
        write_pptx(doc, out)

        prs = PptxPresentation(str(out))
        slide = prs.slides[0]
        assert slide.has_notes_slide
        assert "emphasize X" in slide.notes_slide.notes_text_frame.text

    def test_multi_paragraph_notes(self, tmp_path: Path) -> None:
        doc = ContentDocument(
            metadata=DocumentMetadata(title=""),
            body=(
                Heading(depth=1, children=(Text(value="Slide"),)),
                Div(
                    attr=Attr(classes=("speaker-notes",)),
                    children=(
                        Paragraph(children=(Text(value="Note line 1."),)),
                        Paragraph(children=(Text(value="Note line 2."),)),
                    ),
                ),
            ),
        )
        out = tmp_path / "multinotes.pptx"
        write_pptx(doc, out)

        prs = PptxPresentation(str(out))
        slide = prs.slides[0]
        assert slide.has_notes_slide
        notes = slide.notes_slide.notes_text_frame.text
        assert "Note line 1" in notes
        assert "Note line 2" in notes

    def test_no_notes_on_slide_without_speaker_notes_div(self, tmp_path: Path) -> None:
        doc = ContentDocument(
            metadata=DocumentMetadata(title=""),
            body=(
                Heading(depth=1, children=(Text(value="Slide"),)),
                Paragraph(children=(Text(value="Just content."),)),
            ),
        )
        out = tmp_path / "no_notes.pptx"
        write_pptx(doc, out)

        prs = PptxPresentation(str(out))
        slide = prs.slides[0]
        # has_notes_slide may be False or the text is empty
        if slide.has_notes_slide:
            assert not slide.notes_slide.notes_text_frame.text.strip()
