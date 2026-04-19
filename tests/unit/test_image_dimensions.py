"""Tests for Image.width/height round-trip through DOCX reader + PPTX reader/writer."""

from __future__ import annotations

from pathlib import Path

import pytest
from kaos_content.model.blocks import Figure, Paragraph
from kaos_content.model.document import ContentDocument
from kaos_content.model.inlines import Image
from kaos_content.traversal import find_by_type


class TestEmuToPt:
    """Unit tests for the EMU → points conversion helper."""

    def test_conversion_roundtrip(self) -> None:
        from kaos_office.ooxml.namespace import EMU_PER_POINT, emu_to_pt

        assert emu_to_pt(EMU_PER_POINT) == pytest.approx(1.0)
        # 1 inch = 914400 EMU = 72 pt
        assert emu_to_pt(914400) == pytest.approx(72.0)

    def test_zero(self) -> None:
        from kaos_office.ooxml.namespace import emu_to_pt

        assert emu_to_pt(0) == 0.0

    def test_accepts_int_and_float(self) -> None:
        from kaos_office.ooxml.namespace import emu_to_pt

        assert emu_to_pt(12700) == pytest.approx(1.0)
        assert emu_to_pt(12700.0) == pytest.approx(1.0)


class TestBuilderImageDimensions:
    """``DocumentBuilder.image`` now accepts width / height."""

    def test_builder_passes_dimensions(self) -> None:
        from kaos_content.builders.builder import DocumentBuilder

        builder = DocumentBuilder()
        builder.image("pptx://a.png", alt="a", width=100.0, height=50.0)
        doc = builder.build()
        images = list(find_by_type(doc, Image))
        assert len(images) == 1
        assert images[0].width == 100.0
        assert images[0].height == 50.0

    def test_builder_without_dimensions_unchanged(self) -> None:
        from kaos_content.builders.builder import DocumentBuilder

        builder = DocumentBuilder()
        builder.image("pptx://a.png", alt="a")
        doc = builder.build()
        img = next(iter(find_by_type(doc, Image)))
        assert img.width is None
        assert img.height is None


class TestDocxReaderImageDimensions:
    """DOCX reader must extract wp:extent cx/cy into Image.width/height as points."""

    def test_reader_extracts_extent(self, tmp_path: Path) -> None:
        # Build a minimal DOCX with an image and a known wp:extent, then read it back.
        from kaos_content.model.document import ContentDocument as _CD

        docx_with_image = _fixture_with_image()
        if docx_with_image is None:
            pytest.skip("no DOCX fixture containing an image")

        from kaos_office.docx.reader import parse_docx

        doc = parse_docx(docx_with_image)
        images = list(find_by_type(doc, Image))
        assert images, f"expected at least one Image in {docx_with_image.name}"
        # At least one image should carry non-None dimensions if the source had them.
        dimensioned = [img for img in images if img.width is not None and img.height is not None]
        assert dimensioned, (
            f"expected at least one Image with width/height populated from wp:extent; "
            f"got {[(img.src, img.width, img.height) for img in images]}"
        )
        for img in dimensioned:
            # mypy/ty can't narrow through the comprehension above.
            assert img.width is not None and img.width > 0
            assert img.height is not None and img.height > 0
        # sanity: making sure _CD is actually the kaos-content type
        assert isinstance(doc, _CD)


def _fixture_with_image() -> Path | None:
    """Find a DOCX fixture under kaos-office tests that contains an embedded image."""
    import zipfile

    fixtures_dir = Path(__file__).parent.parent / "fixtures" / "docx"
    if not fixtures_dir.exists():
        return None
    for candidate in sorted(fixtures_dir.glob("*.docx")):
        try:
            with zipfile.ZipFile(candidate) as zf:
                for name in zf.namelist():
                    if name.startswith("word/media/"):
                        return candidate
        except zipfile.BadZipFile:
            continue
    return None


class TestPptxReaderImageDimensions:
    """PPTX reader must read picture-shape EMU dimensions into points."""

    def test_reader_extracts_picture_dimensions(self) -> None:
        pptx_path = Path(__file__).parent.parent / "fixtures" / "pptx" / "battle" / "images.pptx"
        if not pptx_path.exists():
            pytest.skip("no PPTX image fixture")

        from kaos_office.pptx.reader import parse_pptx

        doc = parse_pptx(pptx_path)
        images = list(find_by_type(doc, Image))
        assert images, f"expected at least one Image in {pptx_path.name}"
        with_dims = [img for img in images if img.width is not None and img.height is not None]
        assert with_dims, (
            f"expected at least one Image with width/height populated; "
            f"got {[(img.src, img.width, img.height) for img in images]}"
        )


class TestPptxWriterRespectsDimensions:
    """PPTX writer must pass Image.width/height to ``add_picture`` as EMU."""

    def test_writer_uses_explicit_dimensions(self, tmp_path: Path) -> None:
        pytest.importorskip("pptx")
        from pptx import Presentation

        from kaos_office.pptx.writer import write_pptx

        # 1x1 PNG so we have a real image file
        png_path = tmp_path / "dot.png"
        png_path.write_bytes(_one_pixel_png())

        doc = ContentDocument(
            body=(
                Figure(
                    children=(
                        Paragraph(
                            children=(
                                Image(
                                    src=str(png_path),
                                    alt="dot",
                                    width=144.0,  # 2 inches in points
                                    height=72.0,  # 1 inch in points
                                ),
                            )
                        ),
                    )
                ),
            )
        )
        out = tmp_path / "out.pptx"
        write_pptx(doc, out)
        assert out.exists()

        # Read it back via python-pptx and verify the picture shape's EMU size.
        prs = Presentation(str(out))
        pics = [
            shape
            for slide in prs.slides
            for shape in slide.shapes
            if shape.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
        ]
        assert pics, "expected one picture shape in output"
        # 144 pt = 2 in = 1828800 EMU ; 72 pt = 1 in = 914400 EMU.
        # python-pptx may round to integer EMU; allow a small tolerance.
        assert pics[0].width == pytest.approx(1828800, abs=1)
        assert pics[0].height == pytest.approx(914400, abs=1)


def _one_pixel_png() -> bytes:
    """Return the bytes of a single-pixel opaque white PNG."""
    # Smallest valid PNG — hand-crafted (IHDR 1x1, no compression gotchas).
    return (
        b"\x89PNG\r\n\x1a\n"
        b"\x00\x00\x00\rIHDR"
        b"\x00\x00\x00\x01\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde"
        b"\x00\x00\x00\x0cIDATx\x9cc\xf8\xff\xff?\x00\x05\xfe\x02\xfe"
        b"\xdc\xccY\xe7"
        b"\x00\x00\x00\x00IEND\xaeB`\x82"
    )
